"""
Module for parsing PDF and PPTX files into slides.
"""
import os
from typing import List, Tuple, Dict, Any
import pdfplumber
from pptx import Presentation
from backend.models.core import Slide

def extract_from_file(file_path: str) -> Tuple[List[Slide], Dict[str, Any]]:
    """
    Extracts slides and metadata from a PDF or PPTX file.

    Args:
        file_path: Path to the file.

    Returns:
        Tuple containing:
        - List of Slide objects.
        - Dictionary of deck metadata (author, created_at, etc.)

    Raises:
        ValueError: If the file format is not supported.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        slides, metadata = _extract_from_pdf(file_path)
        # Check for low text density (OCR trigger)
        total_text = sum([len(s.text) for s in slides])
        if total_text < 100 and len(slides) > 0: # Heuristic: <100 chars total for the deck?
             print("Low text density detected. Attempting OCR...")
             ocr_slides = _extract_with_ocr(file_path)
             if ocr_slides:
                 return ocr_slides, metadata
        return slides, metadata

    if ext == ".pptx":
        try:
            return _extract_from_pptx(file_path)
        except Exception as e:
            print(f"PPTX extraction failed: {e}. Attempting PDF conversion fallback (not fully implemented) or OCR if possible.")
            # Fallback: If we could convert PPTX to PDF here, we would. 
            # For now, just re-raise or return empty.
            raise e

    raise ValueError(f"Unsupported file format: {ext}")

def _extract_from_pdf(file_path: str) -> Tuple[List[Slide], Dict[str, Any]]:
    slides = []
    metadata = {}
    with pdfplumber.open(file_path) as pdf:
        metadata = pdf.metadata or {}
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            # PDF doesn't have explicit titles, so we use a generic one or first line
            title = f"Slide {i + 1}"
            slides.append(Slide(index=i, title=title, text=text))
    return slides, metadata

def _extract_from_pptx(file_path: str) -> Tuple[List[Slide], Dict[str, Any]]:
    slides = []
    metadata = {}
    prs = Presentation(file_path)
    
    # Extract metadata
    if prs.core_properties:
        metadata = {
            "author": prs.core_properties.author,
            "created": str(prs.core_properties.created),
            "modified": str(prs.core_properties.modified),
            "title": prs.core_properties.title
        }

    for i, slide in enumerate(prs.slides):
        title = ""
        text_parts = []

        # Try to get title
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text
        else:
            title = f"Slide {i + 1}"

        # Get all text recursively
        for shape in slide.shapes:
            text_parts.append(_get_shape_text(shape))

        # Filter empty strings and join
        text = "\n".join([t for t in text_parts if t.strip()])

        # Get notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text

        slides.append(Slide(index=i, title=title, text=text, notes=notes))

    return slides, metadata

def _get_shape_text(shape: Any) -> str:
    """
    Recursively extracts text from a shape, handling groups and tables.
    """
    text_parts = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text_parts.append("".join(run.text for run in paragraph.runs))

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text_parts.append(_get_shape_text(cell))
    
    # Handle groups
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for sub_shape in shape.shapes:
            text_parts.append(_get_shape_text(sub_shape))

    return "\n".join([t for t in text_parts if t.strip()])

def _extract_with_ocr(file_path: str) -> List[Slide]:
    """
    Extracts text from a file using OCR (converts to images first).
    Useful for image-heavy PDFs or as a fallback.
    """
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except ImportError:
        # If dependencies are missing, return empty or raise
        # For MVP, we'll just log and return empty to avoid crashing if env is partial
        print("OCR dependencies missing (pdf2image or pytesseract). Skipping OCR.")
        return []

    slides = []
    try:
        images = convert_from_path(file_path)
        for i, image in enumerate(images):
            text = pytesseract.image_to_string(image)
            title = f"Slide {i + 1} (OCR)"
            slides.append(Slide(index=i, title=title, text=text))
    except Exception as e:
        print(f"OCR failed: {e}")
        return []

    return slides

def validate_ocr_dependencies():
    """
    Validates that system dependencies for OCR are installed.
    Logs warnings if Tesseract or Poppler are missing.
    """
    import logging
    logger = logging.getLogger(__name__)
    
    # Check Tesseract
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
    except ImportError:
        logger.warning("pytesseract not installed. OCR will not work.")
    except Exception as e:
        logger.warning(f"Tesseract binary not found or not working: {e}. OCR will not work. apt-get install tesseract-ocr / brew install tesseract")

    # Check Poppler (via pdf2image)
    try:
        from pdf2image import pdfinfo_from_path
        # We can't easily check poppler without a file, but we can check if the module is loadable
        # and if the binary is in path. 
        # pdf2image usually checks path on import or usage.
        # Let's try to check visibility of pdfinfo
        from pdf2image.exceptions import PDFInfoNotInstalledError
        try:
             # Basic check if tools are available
             import subprocess
             subprocess.run(["pdfinfo", "-v"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
        except (subprocess.CalledProcessError, FileNotFoundError):
             logger.warning("Poppler ('pdfinfo') not found. PDF conversion/OCR will fail. apt-get install poppler-utils / brew install poppler")
    except ImportError:
        logger.warning("pdf2image not installed. OCR will not work.")
